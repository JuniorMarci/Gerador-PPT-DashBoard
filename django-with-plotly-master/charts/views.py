from django.conf import settings
from django.shortcuts import render
import random
import os
import pandas as pd
import numpy as np
from time import time

def color_cells(val):
    """
    Cores as células com base no valor.
    """
    if val < 90:
        color = 'red'
    elif val < 100:
        color = '#FFD700'
    else:
        color = 'green'
    return 'background-color: %s' % color


def transformar(item):
	try: item = item.replace(',','.')
	except: 1==1
	item = "{:.2f}".format(float(item))
	#print(float(item))
	return float(item)
	

	
	
def GerarHTabela(df,x,y):
	#df = pd.DataFrame(list(zip(df[x], df['Nota Acumulado'])), columns =[x, 'Nota Acumulado'])
	# Aplica a função color_cells para colorir o fundo das células
	#df = df.style.applymap(color_cells, subset=['Nota Pontual','Nota Acumulado'])
	df[x] = [i.split("::")[0] for i in df[x+"m"]]
	df['mês'] = [int(i) for i in df['mês']]
	df = df.sort_values([x,'mês'])
	df = df[[x, 'mês',
	'Previsto Pontual',
	'Realizado Pontual',
	'Nota Pontual',
	'Previsto Acumulado',
	'Realizado Acumulado',
	'Nota Acumulado']]
	
	df = df.style.map(color_cells, subset=['Nota Pontual','Nota Acumulado'])
	html_table = df.to_html(index=False)
	#print(html_table)
	
	return html_table
	
def FormatarTabela(gtab,x,y,MesSel):
	#gtab[x] = 
	#gtab['mêss'] = [str(i) for i in gtab['mês']] # Mês string para concatenar com 'x'
	gtab[x+"m"] = gtab[x]+"::"+[str(i) for i in gtab['mês']]
	mtab = gtab.groupby(x+'m').mean(y).round(2).reset_index().sort_values(y,ascending=False)

	gtab = gtab[gtab['mês'].isin(MesSel)] if len(MesSel) > 0 else gtab

	gtab = gtab.groupby(x).mean(y).reset_index().sort_values(y,ascending=False)
	#tab = tab.dropna()

	vx = gtab[x].values.tolist()
	vy = gtab[y].values.tolist()
	#y = [transformar(i) for i in tab['Realizado Acumulado'].values.tolist()]
	
	return vx,vy,GerarHTabela(mtab,x,y)
	


	
def GerarGrafico(x,y,titulo,conteudo):
	
        

	# Bar chart colors
        #colors = ['limegreen', ] * 5
        colors = ['red' if value < 90 else '#FFD700' if value < 100 else 'green' for value in y]
    	
        # Create a trace json to hold graph data
        trace = {
        'x': x,
        'y': y,
        'type': 'bar',
        'name': 'nome',
        'marker': {'color': colors},
		'text': [round(v,2) for v in y],
		'textposition': 'auto'

        }

        # Configure the chart's layout
        layout = {'title': {'text': titulo+" x "+conteudo,
                        'font': {
                            'color': '#000000'}
                        },
              'xaxis': {'title': titulo, 'color': '#000000', 'mirror': 'true', 'showline': 'false'},
              'yaxis': {'title': conteudo, 'color': '#000000', 'mirror': 'true', 'showline': 'true'},
              'plot_bgcolor': 'white', 'paper_bgcolor': 'white', 'bordercolor': '#000000'}

        # Pass trace and layout in the context
        context = {"trace": trace,
               "layout": layout
               }
			   
        return context
		
def CercareImmagini():
	disp = []
	cam = os.path.join(settings.BASE_DIR,'static/imgs/')
	for c in os.listdir(cam):
		disp.append(c.split("_")[0])
	
	return(set(disp))
	
	

def bar_chart_view(request):

    Disponiveis = CercareImmagini()

    # VARIÁVEIS PARA FILTROS
    sessao = str(time())
    FilDepart = ['9']
    MesSel = [9]
    DepartSel = []
    ShSel = []
    LogSel0 = []
    LogSel1 = []
    LogSel2 = []
	
    # VARIÁVEIS DE CONTEÚDO
    titulo0 = "Departamento"
    titulo1 = "Shopping"
    titulo2 = "Objetivo"
    conteudo0 = "Nota Acumulado"
    conteudo1 = "Nota Acumulado"
    conteudo2 = "Nota Acumulado"
    
	
    if request.method == 'POST':
        FilDepart = []
        for key in request.POST:
            if key == 'csrfmiddlewaretoken':
                continue
            #print(f'{key}: {request.POST.get(key)}')
            #FilDepart.append(f'{request.POST.get(key)}')
            MesSel = [int(i) for i in request.POST.getlist('mes')]
            DepartSel = request.POST.getlist('departamento')
            ShSel = request.POST.getlist('shopping')
            sessao = request.POST.getlist('sessao')[0]
            titulo0 = request.POST.getlist('titulo0')[0]
            titulo1 = request.POST.getlist('titulo1')[0]
            titulo2 = request.POST.getlist('titulo2')[0]
            LogSel0 = request.POST.getlist('login')

	# TRATANDO A TABELA PRINCIPAL
    FilDepart = [int(i) if i.isdigit() else i for i in FilDepart]
    tab = os.path.join(settings.BASE_DIR, 'static', 'Metas.csv')
    tab = pd.read_csv(tab,sep=';')
    tab['Departamento'] = [c.split('_')[0] for c in tab['Código da Meta']]
	
    # LISTA DE FILTROS FIXA
    Depart = sorted(list(set(tab['Departamento'])))#[c.split('_')[0] for c in tab['Código da Meta']])
    Meses = set([int(i) for i in tab['mês'].values.tolist()])
    #tab = tab[tab['mês']==9]
    tab = tab.rename(columns={"Código da Meta": "Código_da_Meta"})
	
	# INCLUINDO COLUNA DE SHOPPINGS
    Shoppings = os.path.join(settings.BASE_DIR, 'static', 'Categorias.csv')
    Shoppings = pd.read_csv(Shoppings,sep=';')
    tab = tab.merge(Shoppings,left_on='Cod.Unico',right_on='0',how='left').drop(columns=["0"]).rename(columns={'1':'Shopping'})
    sh = sorted(list(set(tab['Shopping'])))
	
	
    # FILTRANDO CONFORME O SOLICITADO PELO USUÁRIO    
    #print(MesSel, DepartSel, ShSel)
    tab = tab[tab['Departamento'].isin(DepartSel)] if len(DepartSel) > 0 else tab
    #tab = tab[tab['mês'].isin(MesSel)] if len(MesSel) > 0 else tab
    tab = tab[tab['Shopping'].isin(ShSel)] if len(ShSel) > 0 else tab
    tab = tab[tab['Login'].isin(LogSel0)] if len(LogSel0) > 0 else tab
	
    LogDisp = set(tab['Login'])
	
    def TransNum(tab):
	    # TRATANDO OS NÚMEROS
	    tab['Nota Acumulado'] = [transformar(i) for i in tab['Nota Acumulado'].values.tolist()]
	    tab = tab[tab['Nota Acumulado']!=-1]
	    tab = tab[~tab['Nota Acumulado'].isna()]
	    tab['Nota Pontual'] = [transformar(i) for i in tab['Nota Pontual'].values.tolist()]
	    tab = tab[tab['Nota Pontual']!=-1]
	    tab = tab[~tab['Nota Pontual'].isna()]
	    tab['Realizado Pontual'] = [transformar(i) for i in tab['Realizado Pontual'].values.tolist()]
	    tab['Realizado Acumulado'] = [transformar(i) for i in tab['Realizado Acumulado'].values.tolist()]
	    tab['Previsto Pontual'] = [transformar(i) for i in tab['Previsto Pontual'].values.tolist()]
	    tab['Previsto Acumulado'] = [transformar(i) for i in tab['Previsto Acumulado'].values.tolist()]
	    return tab

    tab = TransNum(tab)
		
    x0,y0,t0 = FormatarTabela(tab,titulo0,conteudo0,MesSel)
    x1,y1,t1 = FormatarTabela(tab,titulo1,conteudo1,MesSel)
    x2,y2,t2 = FormatarTabela(tab,titulo2,conteudo2,MesSel)
	
    
    
		
    G0 = GerarGrafico(x0,y0,titulo0,conteudo0)
    G0['table'] = t0
    G1 = GerarGrafico(x1,y1,titulo1,conteudo1)
    G1['table'] = t1
    G2 = GerarGrafico(x2,y2,titulo2,conteudo2)
    G2['table'] = t2
	
    context = {"title": "Gerador de Apresentações 2.0",
				"FilDepart":FilDepart,
				"Depart":Depart, # Lista da tabela
				"Meses":Meses, # Lista da tabela
				"LogDisp":LogDisp, # Lista da tabela (pós filtragem)
				"Titulo":['Login','Código_da_Meta','Objetivo','Departamento','Shopping'],
				"MesSel":MesSel, # Origem: HTML
				"DepartSel":DepartSel, # Origem: HTML
				"ShSel":ShSel, # Origem: HTML
				"LogSel0":LogSel0,# Origem: HTML
				"LogSel1":LogSel1,# Origem: HTML
				"LogSel2":LogSel2,# Origem: HTML
				"sessao":sessao, # Origem: HTML
				"sh":sh,
				'G0':G0,
				'G1':G1,
				'G2':G2,
				'titulo0':titulo0, # Origem: HTML
				'titulo1':titulo1, # Origem: HTML
				'titulo2':titulo2, # Origem: HTML
				'Disponiveis':Disponiveis

			   }
	
    #print(context)	
    #print(tab.columns)

    return render(request, "chart.html", context)


def scatter_chart_view(request):
    """
    Scatter Chart
    :param request:
    :return:
    """
    # Generate a list of random numbers
    x = list(range(1, 6))
    y = [random.randint(100, 600) for _ in range(6)]

    # Create a trace json to hold graph data
    trace = {
        'x': x,
        'y': y,
        'type': 'scatter',

    }

    # Configure the chart's layout
    layout = {'title': {'text': 'Scatter Chart',
                        'font': {
                            'color': '#ffffff'}
                        },
              'xaxis': {'title': 'X-axis', 'color': '#DCDCDC', 'mirror': 'true', 'showline': 'false'},
              'yaxis': {'title': 'Y-axis', 'color': '#DCDCDC', 'mirror': 'true', 'showline': 'true'},
              'plot_bgcolor': 'black', 'paper_bgcolor': 'black', 'bordercolor': '#ffffff'}

    # Pass trace and layout in the context
    context = {"trace": trace,
               "layout": layout,
               "title": "Scatter Chart"}

    return render(request, "chart.html", context)


def pie_chart_view(request):
    """
    Pie Chart
    :param request:
    :return:
    """

    # Create a trace json to hold graph data
    trace = {
        'values': [19, 26, 55],
        'labels': ['Residential', 'Non-Residential', 'Utility'],
        'type': 'pie'


    }

    # Configure the chart's layout
    layout = {'title': {'text': 'Pie Chart',
                        'font': {
                            'color': '#ffffff'}
                        },
              'plot_bgcolor': 'black', 'paper_bgcolor': 'black', 'bordercolor': '#ffffff'}

    # Pass trace and layout in the context
    context = {"trace": trace,
               "layout": layout,
               "title": "Pie Chart"}

    return render(request, "chart.html", context)


def bubble_chart_view(request):
    """
    Bubble Chart
    :param request:
    :return:
    """

    # Create a trace json to hold graph data
    trace = {
          'x': [1, 2, 3, 4],
          'y': [10, 11, 12, 13],
          'mode': 'markers',
          'marker': {
            'size': [40, 60, 80, 100]
          }


    }

    # Configure the chart's layout
    layout = {'title': {'text': 'Bubble Chart',
                        'font': {
                            'color': '#ffffff'}
                        },
              'xaxis': {'title': 'X-axis', 'color': '#DCDCDC', 'mirror': 'true', 'showline': 'false'},
              'yaxis': {'title': 'Y-axis', 'color': '#DCDCDC', 'mirror': 'true', 'showline': 'true'},
              'plot_bgcolor': 'black', 'paper_bgcolor': 'black', 'bordercolor': '#ffffff'}

    # Pass trace and layout in the context
    context = {"trace": trace,
               "layout": layout,
               "title": "Bubble Chart"
			   }

    return render(request, "chart.html", context)
	
	
	
	

	
	
	
from django.http import JsonResponse
from django.views.decorators.csrf import csrf_exempt

@csrf_exempt
def upload_image(request,sessao,nome):
    if request.method == 'POST':
        file = request.FILES['file']
        handle_uploaded_file(file,sessao,nome)
        return JsonResponse({'status': 'success'})
    else:
        return JsonResponse({'status': 'error'})

		
def handle_uploaded_file(file,sessao,nome):
    file_path = os.path.join(settings.BASE_DIR, 'static/imgs', sessao+'_'+nome+'.png')
    with open(file_path, 'wb+') as destination:
        for chunk in file.chunks():
            destination.write(chunk)
			
			

def GerarApresentacao(response,sessao):
			
	from pptx import Presentation
	from pptx.util import Inches
	from django.http import FileResponse

	# Crie uma instância da classe Presentation
	prs = Presentation()

	# Caminho para a imagem
	img_path = os.path.join(settings.BASE_DIR, 'static/imgs/')
	for c in os.listdir(img_path):
		if (c.split("_")[0] == sessao):
			img = os.path.join(settings.BASE_DIR, 'static/imgs/'+c)

			# Adicione um slide com um layout em branco
			blank_slide_layout = prs.slide_layouts[6]
			slide = prs.slides.add_slide(blank_slide_layout)

			# Defina a posição da imagem (aqui é 1 polegada da borda superior e esquerda)
			left = top = Inches(1)

			# Adicione a imagem ao slide
			pic = slide.shapes.add_picture(img, left, top)

	# Salve a apresentação
	pptx_path = os.path.join(settings.BASE_DIR,'static/apresentacoes/'+sessao+'.pptx')
	prs.save(pptx_path)
	
	return FileResponse(open(pptx_path, 'rb'), as_attachment=True, filename=sessao+'.pptx')